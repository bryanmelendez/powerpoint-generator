from pptx import Presentation
from pptx.util import Inches


class Slide():

    # make class attributes stuff like slide dimensions, font sizes, etc

    def __init__(self, title_text):
        self.title_text = title_text


class TitleSlide(Slide):
    def __init__(self, title_text, subtitle_text):
        self.subtitle_text = subtitle_text
        super().__init__(title_text)

    def CreateSlide(self, pres):
        title_slide_layout = pres.slide_layouts[0]
        slide = pres.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = self.title_text
        subtitle.text = self.subtitle_text


class ImageSlide(Slide):
    def __init__(self, title_text, image_path):
        self.image_path = image_path
        super().__init__(title_text)

    def CreateSlide(self, pres):
        image_slide_layout = pres.slide_layouts[6]
        slide = pres.slides.add_slide(image_slide_layout)

        # Figure out what layout i need to make it have a title
        # title = slide.shapes.title
        # title.text = self.title_text

        left = Inches(1)
        top = Inches(1)
        pic = slide.shapes.add_picture(self.image_path, left, top)
       
# Driver code
pres = Presentation()

title_slide = TitleSlide('Example title', 'Example subtitle')
title_slide.CreateSlide(pres)

image_slide = ImageSlide('Image slide title', '{}Example Exports/Option 1 A.jpg'.format(EXAMPLE_IMAGES))
image_slide.CreateSlide(pres)

pres.save('test.pptx')
