# Classes for making pptx slides
from pptx import Presentation
from pptx.util import Inches
from PIL import Image


class Slide():
    SLD_LAYOUT_TITLE_SLIDE = 0
    SLD_LAYOUT_TITLE_AND_CONTENT = 1
    SLD_LAYOUT_PICTURE_WITH_CAPTION = 8

    PICTURE_PLACEHOLDER_IDX = 13
    TEXT_PLACEHOLDER_IDX = 14

    def __init__(self, title_text):
        self.title_text = title_text


class TitleSlide(Slide):
    def __init__(self, title_text, subtitle_text):
        self.subtitle_text = subtitle_text
        super().__init__(title_text)

    def create_slide(self, pres):
        title_slide_layout = pres.slide_layouts[self.SLD_LAYOUT_TITLE_SLIDE]
        slide = pres.slides.add_slide(title_slide_layout)

        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = self.title_text
        subtitle.text = self.subtitle_text


class ImageSlide(Slide):
    def __init__(self, title_text, footer_text, image_path):
        self.image_path = image_path
        self.footer_text = footer_text
        super().__init__(title_text)

    def create_slide(self, pres):
        image_slide_layout = pres.slide_layouts[self.SLD_LAYOUT_TITLE_AND_CONTENT]
        slide = pres.slides.add_slide(image_slide_layout)

        # Debug template

        # for shape in slide.placeholders:
        #     print('%d %s' % (shape.placeholder_format.idx, shape.name))

        # quit()

        pic_placeholder = slide.placeholders[self.PICTURE_PLACEHOLDER_IDX]

        picture = pic_placeholder.insert_picture(self.image_path)

        self.resize_image(picture)

        title_placeholder = slide.shapes.title
        title_placeholder.text = self.title_text

        text_placeholder = slide.placeholders[self.TEXT_PLACEHOLDER_IDX]
        text_placeholder.text = self.footer_text

    # fits the image inside the existing placeholder
    # credits to: https://stackoverflow.com/questions/56815178/how-can-i-get-the-dimensions-of-a-picture-placeholder-to-re-size-an-image-when-c
    def resize_image(self, picture):
        picture.crop_top = 0
        picture.crop_left = 0
        picture.crop_bottom = 0
        picture.crop_right = 0

        width, height = picture.image.size  # ---width and height are int pixel-counts

        available_width = picture.width
        available_height = picture.height
        image_width, image_height = picture.image.size
        placeholder_aspect_ratio = float(available_width) / float(available_height)
        image_aspect_ratio = float(image_width) / float(image_height)

        # Get initial image placeholder left and top positions
        pos_left, pos_top = picture.left, picture.top

        picture.crop_top = 0
        picture.crop_left = 0
        picture.crop_bottom = 0
        picture.crop_right = 0

        # ---if the placeholder is "wider" in aspect, shrink the picture width while
        # ---maintaining the image aspect ratio
        if placeholder_aspect_ratio > image_aspect_ratio:
            picture.width = int(image_aspect_ratio * available_height)
            picture.height = available_height

        # ---otherwise shrink the height
        else:
            picture.height = int(available_width/image_aspect_ratio)
            picture.width = available_width

        # Set the picture left and top position to the initial placeholder one
        picture.left = pos_left
        picture.top = pos_top
